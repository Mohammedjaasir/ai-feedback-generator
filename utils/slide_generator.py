from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_bulleted_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    textbox = slide.shapes.placeholders[1]
    text_frame = textbox.text_frame
    text_frame.clear()

    for line in content.split('\n'):
        line = line.strip()
        if line:
            p = text_frame.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = 'Calibri'
            p.font.color.rgb = RGBColor(40, 40, 40)

def create_feedback_slides(feedback: str, improvements: str, output_file: str = "Project_Feedback_Slides.pptx"):
    prs = Presentation()

    # Slide 1 - Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Feedback Report"
    subtitle = slide.placeholders[1]
    subtitle.text = "Auto-generated with Groq + LLMs"
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)

    # Slide 2 - Feedback Summary
    add_bulleted_slide(prs, "📝 Feedback Summary", feedback)

    # Slide 3 - Improvement Ideas
    add_bulleted_slide(prs, "💡 Project Improvement Ideas", improvements)

    # Slide 4 (Optional) - Viva Questions (auto-detected)
    if "viva" in feedback.lower() or "questions" in feedback.lower():
        lines = [line for line in feedback.split('\n') if "?" in line or "question" in line.lower()]
        if lines:
            add_bulleted_slide(prs, "🎯 Possible Viva Questions", "\n".join(lines))

    prs.save(output_file)
    return output_file
